#!/usr/bin/env python3
"""WAWA SmartERP 사용자가이드 — 누락 기능 슬라이드 추가"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from copy import deepcopy
from lxml import etree
import os

INPUT  = "WAWA_SmartERP_사용자가이드.pptx"
OUTPUT = "WAWA_SmartERP_사용자가이드.pptx"
SS     = "apps/desktop/e2e-screenshots-all"

# ── 기존 가이드 색상 시스템 (역공학) ─────────────────────
NAVY      = RGBColor(0x0F, 0x17, 0x2A)   # 메인 배경
BLUE_HDR  = RGBColor(0x33, 0x72, 0xF7)   # 헤더 바
BG_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BG_BOTTOM = RGBColor(0x1E, 0x3A, 0x5F)   # 하단 진한 블루
TEXT_WH   = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_SUB  = RGBColor(0xBF, 0xD7, 0xFF)   # 연한 부제목
TEXT_DK   = RGBColor(0x1E, 0x29, 0x3B)   # 다크 본문
TEXT_GY   = RGBColor(0x64, 0x74, 0x8B)   # 회색 설명
ACCENT_G  = RGBColor(0x16, 0xA3, 0x4A)   # 초록
ACCENT_O  = RGBColor(0xEA, 0x58, 0x0C)   # 주황
KAKAO_Y   = RGBColor(0xFF, 0xE8, 0x12)   # 카카오

W = Inches(13.33)
H = Inches(7.5)


# ── 유틸 ─────────────────────────────────────────────────
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, x, y, w, h, fill=None, border=None, border_w=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.line.fill.background()
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border
        if border_w: s.line.width = border_w
    else:
        s.line.fill.background()
    return s


def txt(slide, text, x, y, w, h, size=14, bold=False,
        color=TEXT_DK, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return box


def img_fit(slide, path, x, y, w, h):
    if not os.path.exists(path):
        rect(slide, x, y, w, h, fill=RGBColor(0xDD, 0xDD, 0xDD))
        txt(slide, f"[없음]\n{os.path.basename(path)}", x, y, w, h,
            size=9, color=TEXT_GY, align=PP_ALIGN.CENTER)
        return
    from PIL import Image
    iw, ih = Image.open(path).size
    ratio  = min(w / iw / 9144, h / ih / 9144)
    nw, nh = int(iw * 9144 * ratio), int(ih * 9144 * ratio)
    ox, oy = x + (w - nw)//2, y + (h - nh)//2
    slide.shapes.add_picture(path, ox, oy, nw, nh)


def header(slide, title, subtitle="", step=None, step_color=BLUE_HDR):
    """기존 가이드와 동일한 헤더 스타일"""
    rect(slide, 0, 0, W, Inches(1.2), fill=NAVY)
    txt(slide, title, Inches(0.45), Inches(0.1), Inches(9.5), Inches(0.72),
        size=28, bold=True, color=TEXT_WH)
    if subtitle:
        txt(slide, subtitle, Inches(0.45), Inches(0.77), Inches(10), Inches(0.38),
            size=14, color=TEXT_SUB)
    if step:
        rect(slide, Inches(10.8), Inches(0.22), Inches(2.1), Inches(0.72), fill=step_color)
        txt(slide, step, Inches(10.8), Inches(0.22), Inches(2.1), Inches(0.72),
            size=13, bold=True, color=NAVY if step_color == KAKAO_Y else TEXT_WH,
            align=PP_ALIGN.CENTER)


def tip_box(slide, items, x=Inches(8.2), y=Inches(2.4),
            w=Inches(4.7), title="💡  이렇게 하세요"):
    """우측 팁 박스 — 기존 가이드 스타일"""
    rect(slide, x, y, w, H - y - Inches(0.3),
         fill=RGBColor(0x1E, 0x3A, 0x5F))
    txt(slide, title, x + Inches(0.2), y + Inches(0.1),
        w - Inches(0.4), Inches(0.45),
        size=13, bold=True, color=TEXT_WH)
    for i, (icon, head, body) in enumerate(items):
        oy = y + Inches(0.7 + i * 1.35)
        txt(slide, f"{icon}  {head}", x + Inches(0.2), oy,
            w - Inches(0.4), Inches(0.38),
            size=12, bold=True, color=KAKAO_Y)
        txt(slide, body, x + Inches(0.2), oy + Inches(0.42),
            w - Inches(0.4), Inches(0.82),
            size=11, color=TEXT_SUB)


def step_badge(slide, n, label, color, x, y):
    rect(slide, x, y, Inches(0.5), Inches(0.5), fill=color)
    txt(slide, str(n), x, y, Inches(0.5), Inches(0.5),
        size=14, bold=True, color=TEXT_WH, align=PP_ALIGN.CENTER)
    txt(slide, label, x + Inches(0.6), y + Inches(0.05),
        Inches(4), Inches(0.4), size=13, bold=True, color=TEXT_DK)


# ════════════════════════════════════════════════════════
#  새 슬라이드 정의
# ════════════════════════════════════════════════════════

def slide_timer_intro(prs):
    """타이머 기능 소개 섹션 구분"""
    s = blank(prs)
    bg(s, NAVY)
    rect(s, 0, Inches(3.2), W, Inches(0.07), fill=BLUE_HDR)
    txt(s, "⏱  수업 타이머 사용법",
        Inches(1), Inches(1.8), Inches(11), Inches(1.2),
        size=40, bold=True, color=TEXT_WH, align=PP_ALIGN.CENTER)
    txt(s, "학생 입장 체크인 · 수업 연장 · 일시정지 · 퇴원까지",
        Inches(1), Inches(3.4), Inches(11), Inches(0.7),
        size=18, color=TEXT_SUB, align=PP_ALIGN.CENTER)
    txt(s, "상단 메뉴 → 타이머 탭 선택",
        Inches(1), Inches(4.3), Inches(11), Inches(0.5),
        size=14, color=RGBColor(0x80, 0x9B, 0xBD), align=PP_ALIGN.CENTER)


def slide_timer_checkin(prs):
    """⑦ 타이머 — 대기 학생 & 체크인"""
    s = blank(prs)
    bg(s, BG_WHITE)
    header(s, "⑦ 타이머 — 학생 체크인",
           "타이머 탭 > 실시간 현황 — 대기 학생을 클릭하면 수업이 시작됩니다",
           step="STEP 7", step_color=BLUE_HDR)

    # 왼쪽: 대기 목록 스크린샷
    rect(s, Inches(0.45), Inches(1.35), Inches(5.5), Inches(5.0),
         fill=RGBColor(0xF1,0xF5,0xF9))
    img_fit(s, f"{SS}/rt_03_대기_학생_목록.png",
            Inches(0.5), Inches(1.4), Inches(5.4), Inches(4.8))
    txt(s, "① 대기 중인 학생 목록", Inches(0.5), Inches(6.3), Inches(5.4), Inches(0.38),
        size=11, bold=True, color=TEXT_GY, align=PP_ALIGN.CENTER)

    # 오른쪽: 체크인 후 타이머 카드
    rect(s, Inches(6.2), Inches(1.35), Inches(6.5), Inches(5.0),
         fill=RGBColor(0xF1,0xF5,0xF9))
    img_fit(s, f"{SS}/rt_04_체크인_수업카드.png",
            Inches(6.25), Inches(1.4), Inches(6.4), Inches(4.8))
    txt(s, "② 클릭 → 타이머 시작! 남은 시간 실시간 표시",
        Inches(6.25), Inches(6.3), Inches(6.4), Inches(0.38),
        size=11, bold=True, color=TEXT_GY, align=PP_ALIGN.CENTER)

    # 하단 안내
    rect(s, Inches(0.45), Inches(6.75), Inches(12.4), Inches(0.55),
         fill=RGBColor(0xEF,0xF3,0xFF))
    txt(s, "▶  학생 이름 클릭 한 번으로 체크인 완료 — 타이머가 즉시 시작됩니다",
        Inches(0.65), Inches(6.8), Inches(12), Inches(0.45),
        size=12, color=BLUE_HDR)


def slide_timer_pause_extend(prs):
    """⑧ 타이머 — 일시정지 & 수업 연장"""
    s = blank(prs)
    bg(s, BG_WHITE)
    header(s, "⑧ 타이머 — 일시정지 & 수업 연장",
           "수업 중 잠깐 자리를 비우거나, 수업 시간을 늘릴 때 사용합니다",
           step="STEP 8", step_color=ACCENT_O)

    # 일시정지
    rect(s, Inches(0.45), Inches(1.35), Inches(5.9), Inches(5.0),
         fill=RGBColor(0xFF, 0xF7, 0xED))
    img_fit(s, f"{SS}/rt_05_일시정지_사유선택.png",
            Inches(0.5), Inches(1.4), Inches(5.8), Inches(4.0))

    rect(s, Inches(0.45), Inches(5.45), Inches(5.9), Inches(0.55), fill=ACCENT_O)
    txt(s, "일시정지 — 사유 선택 (외출·휴식·화장실·기타)",
        Inches(0.55), Inches(5.48), Inches(5.7), Inches(0.48),
        size=12, bold=True, color=TEXT_WH)

    # 수업 연장
    rect(s, Inches(6.55), Inches(1.35), Inches(6.3), Inches(5.0),
         fill=RGBColor(0xF0, 0xFD, 0xF4))
    img_fit(s, f"{SS}/uc05-extend-sheet.png",
            Inches(6.6), Inches(1.4), Inches(6.2), Inches(4.0))

    rect(s, Inches(6.55), Inches(5.45), Inches(6.3), Inches(0.55), fill=ACCENT_G)
    txt(s, "수업 연장 — +10 / +20 / +30분 또는 직접 입력",
        Inches(6.65), Inches(5.48), Inches(6.1), Inches(0.48),
        size=12, bold=True, color=TEXT_WH)

    # 하단 안내
    rect(s, Inches(0.45), Inches(6.75), Inches(12.4), Inches(0.55),
         fill=RGBColor(0xFF, 0xF7, 0xED))
    txt(s, "▶  수업 카드의 '정지' 버튼 → 사유 선택  |  '수업추가' 버튼 → 연장 시간 선택",
        Inches(0.65), Inches(6.8), Inches(12), Inches(0.45),
        size=12, color=ACCENT_O)


def slide_timer_temp(prs):
    """⑨ 타이머 — 임시 학생 추가"""
    s = blank(prs)
    bg(s, BG_WHITE)
    header(s, "⑨ 타이머 — 목록에 없는 학생 즉석 추가",
           "명단에 없는 학생이 갑자기 왔을 때 — 임시 학생으로 바로 등록해 수업 진행",
           step="STEP 9", step_color=BLUE_HDR)

    img_fit(s, f"{SS}/uc03-temp-modal.png",
            Inches(1.5), Inches(1.4), Inches(6.5), Inches(5.0))

    # 우측 설명 카드
    tip_box(s, [
        ("①", "왼쪽 상단 '임시' 버튼 클릭",
         "대기 목록 상단에 있는\n'임시 학생 추가' 버튼을 누르세요"),
        ("②", "이름·학년·시간·과목 입력",
         "수업 시작/종료 시간을\n직접 입력합니다"),
        ("③", "'추가하기' 버튼 클릭",
         "대기 목록에 즉시 나타납니다\n이후 일반 학생과 동일하게 체크인"),
    ], title="임시 학생 추가 순서")

    rect(s, Inches(0.45), Inches(6.75), Inches(12.4), Inches(0.55),
         fill=RGBColor(0xEF, 0xF3, 0xFF))
    txt(s, "▶  임시 학생은 앱 종료 시 자동 삭제됩니다 — 정식 등록은 관리자에게 문의",
        Inches(0.65), Inches(6.8), Inches(12), Inches(0.45),
        size=12, color=BLUE_HDR)


def slide_report_intro(prs):
    """AI 보고서 섹션 구분"""
    s = blank(prs)
    bg(s, NAVY)
    rect(s, 0, Inches(3.2), W, Inches(0.07), fill=ACCENT_O)
    txt(s, "📋  AI 월말 보고서 사용법",
        Inches(1), Inches(1.8), Inches(11), Inches(1.2),
        size=40, bold=True, color=TEXT_WH, align=PP_ALIGN.CENTER)
    txt(s, "점수 입력 → AI 코멘트 자동 생성 → 카카오 발송까지",
        Inches(1), Inches(3.4), Inches(11), Inches(0.7),
        size=18, color=TEXT_SUB, align=PP_ALIGN.CENTER)
    txt(s, "상단 메뉴 → 성적표 탭 선택 (이미 앞에서 실습한 메뉴입니다)",
        Inches(1), Inches(4.3), Inches(11), Inches(0.5),
        size=14, color=RGBColor(0x80, 0x9B, 0xBD), align=PP_ALIGN.CENTER)


def slide_ai_comment(prs):
    """⑩ AI 코멘트 자동 생성"""
    s = blank(prs)
    bg(s, BG_WHITE)
    header(s, "⑩ 성적 입력 — AI 코멘트 자동 생성",
           "점수 저장 후 스크롤을 내리면 AI 코멘트 섹션이 나옵니다",
           step="STEP 10", step_color=ACCENT_O)

    img_fit(s, f"{SS}/13_성적표_AI_설정.png",
            Inches(0.45), Inches(1.35), Inches(7.7), Inches(5.1))

    tip_box(s, [
        ("①", "AI 제공사 선택",
         "Gemini (무료) · ChatGPT · Claude\n중 하나를 선택하세요\n기본값은 Gemini입니다"),
        ("②", "'AI 생성' 버튼 클릭",
         "약 10~30초 후 코멘트가\n자동으로 생성됩니다"),
        ("③", "마음에 드는 버전 선택",
         "버전 1·2·3 중 고른 뒤\n'이 버전 사용' 클릭\n수정도 가능합니다"),
        ("④", "반드시 저장!",
         "코멘트 확인 후\n'저장' 버튼을 눌러야\n미리보기에 반영됩니다"),
    ], title="AI 코멘트 생성 순서")

    rect(s, Inches(0.45), Inches(6.75), Inches(12.4), Inches(0.55),
         fill=RGBColor(0xFF, 0xF7, 0xED))
    txt(s, "▶  AI API 키는 관리자가 설정해 둡니다 — 선생님은 버튼만 누르면 됩니다",
        Inches(0.65), Inches(6.8), Inches(12), Inches(0.45),
        size=12, color=ACCENT_O)


def slide_preview(prs):
    """⑪ 성적표 미리보기"""
    s = blank(prs)
    bg(s, BG_WHITE)
    header(s, "⑪ 성적표 미리보기",
           "성적 입력 > 미리보기 탭 — 학부모에게 발송 전 최종 확인",
           step="STEP 11", step_color=ACCENT_G)

    img_fit(s, f"{SS}/11_성적표_정지효_월말평가서.png",
            Inches(0.45), Inches(1.35), Inches(7.7), Inches(5.1))

    tip_box(s, [
        ("📈", "6개월 성적 트렌드",
         "자동으로 최근 6개월\n성적 변화 그래프가\n생성됩니다"),
        ("📝", "과목별 코멘트",
         "각 선생님이 입력한\n점수와 코멘트가\n과목 카드로 표시됩니다"),
        ("🌟", "종합 평가",
         "AI가 생성하거나 직접 쓴\n종합 코멘트가 노란\n박스로 강조됩니다"),
        ("🖨", "저장 방법",
         "오른쪽 상단 버튼으로\nJPG 저장 또는 인쇄\n가능합니다"),
    ], title="미리보기 구성")

    rect(s, Inches(0.45), Inches(6.75), Inches(12.4), Inches(0.55),
         fill=RGBColor(0xF0, 0xFD, 0xF4))
    txt(s, "▶  발송 전 반드시 미리보기로 내용을 확인하세요 — 이름·점수 오타 최종 검토!",
        Inches(0.65), Inches(6.8), Inches(12), Inches(0.45),
        size=12, color=ACCENT_G)


def slide_kakao_send(prs):
    """⑫ 카카오 알림톡 발송"""
    s = blank(prs)
    bg(s, BG_WHITE)
    header(s, "⑫ 카카오 알림톡 — 학부모 발송",
           "성적 입력 > 발송 탭 — 전체 학생을 한 번에 학부모에게 전송",
           step="STEP 12", step_color=KAKAO_Y)

    img_fit(s, f"{SS}/nav_08_step3_send.png",
            Inches(0.45), Inches(1.35), Inches(7.7), Inches(5.1))

    tip_box(s, [
        ("①", "발송 탭으로 이동",
         "성적표 메뉴 3단계 중\n마지막 '발송' 탭을\n클릭하세요"),
        ("②", "학생 체크박스 선택",
         "전체 선택 또는\n개별 선택 가능합니다"),
        ("③", "'일괄 발송' 클릭",
         "카카오 로그인 상태에서\n선택한 학생 전체에게\n한 번에 발송됩니다"),
        ("④", "발송 완료 확인",
         "상태가 '전송완료'로\n바뀌면 학부모 카카오에\n메시지가 도착합니다"),
    ], title="발송 순서")

    rect(s, Inches(0.45), Inches(6.75), Inches(12.4), Inches(0.55),
         fill=RGBColor(0xFE, 0xFC, 0xE8))
    txt(s, "▶  카카오 로그인이 필요합니다 — 앱 우측 상단 카카오 버튼으로 로그인 후 발송하세요",
        Inches(0.65), Inches(6.8), Inches(12), Inches(0.45),
        size=12, color=RGBColor(0x92, 0x70, 0x00))


def slide_final(prs):
    """완료 슬라이드 (기존 것을 대체)"""
    s = blank(prs)
    bg(s, NAVY)
    rect(s, 0, Inches(4.5), W, H - Inches(4.5), fill=BG_BOTTOM)

    txt(s, "✅  전체 기능 완료!",
        Inches(0.8), Inches(0.65), Inches(11), Inches(1.1),
        size=42, bold=True, color=TEXT_WH)
    txt(s, "이제 Wawa Smart ERP의 모든 핵심 기능을 사용할 수 있습니다.",
        Inches(0.8), Inches(1.8), Inches(11), Inches(0.6),
        size=16, color=TEXT_SUB)
    rect(s, Inches(0.8), Inches(2.55), Inches(7.0), Inches(0.06),
         fill=BLUE_HDR)

    # 기능 목록
    features = [
        ("⑦", "타이머 — 체크인/정지/연장"),
        ("⑧", "일시정지 & 수업 연장"),
        ("⑨", "임시 학생 즉석 추가"),
        ("⑩", "AI 코멘트 자동 생성"),
        ("⑪", "성적표 미리보기"),
        ("⑫", "카카오 알림톡 발송"),
    ]
    for i, (num, label) in enumerate(features):
        col, row = divmod(i, 3)
        x = Inches(0.8 + col * 4.0)
        y = Inches(2.75 + row * 0.75)
        rect(s, x, y + Inches(0.05), Inches(0.4), Inches(0.4),
             fill=BLUE_HDR)
        txt(s, num, x, y + Inches(0.05), Inches(0.4), Inches(0.4),
            size=11, bold=True, color=TEXT_WH, align=PP_ALIGN.CENTER)
        txt(s, label, x + Inches(0.5), y + Inches(0.08),
            Inches(3.3), Inches(0.38), size=12, color=TEXT_WH)

    # 우측 팁
    rect(s, Inches(8.8), Inches(2.55), Inches(4.1), Inches(4.7),
         fill=RGBColor(0x1A, 0x2E, 0x44))
    txt(s, "꼭 기억하세요!", Inches(9.0), Inches(2.7),
        Inches(3.7), Inches(0.45), size=13, bold=True, color=KAKAO_Y)
    tips = [
        ("⚠️  저장 필수", "과목별 입력 후 반드시\n'저장' 버튼 클릭"),
        ("🔑  카카오 로그인", "발송 전 카카오 로그인\n상태 확인 필요"),
        ("👤  학생 추가", "명단 외 학생은 관리자에게\n정식 등록 요청"),
    ]
    for i, (head, body) in enumerate(tips):
        oy = Inches(3.3 + i * 1.3)
        txt(s, head, Inches(9.0), oy, Inches(3.7), Inches(0.38),
            size=12, bold=True, color=KAKAO_Y)
        txt(s, body, Inches(9.0), oy + Inches(0.42), Inches(3.7), Inches(0.75),
            size=11, color=TEXT_SUB)


# ════════════════════════════════════════════════════════
#  MAIN — 기존 PPT 로드 후 마지막 슬라이드 교체 + 새 슬라이드 추가
# ════════════════════════════════════════════════════════

def main():
    prs = Presentation(INPUT)

    # 기존 마지막 슬라이드(완료 슬라이드) 삭제 → 나중에 새걸로 교체
    # pptx는 슬라이드 삭제 API가 없으므로 XML에서 직접 제거
    xml_slides = prs.slides._sldIdLst
    last_ref   = xml_slides[-1]
    xml_slides.remove(last_ref)

    # 새 슬라이드 추가 (순서대로)
    slide_timer_intro(prs)
    slide_timer_checkin(prs)
    slide_timer_pause_extend(prs)
    slide_timer_temp(prs)
    slide_report_intro(prs)
    slide_ai_comment(prs)
    slide_preview(prs)
    slide_kakao_send(prs)
    slide_final(prs)

    prs.save(OUTPUT)
    print(f"✅  저장 완료: {OUTPUT}")
    print(f"   총 슬라이드: {len(prs.slides)}장")
    print()
    print("  기존 (1~13)  — 다운로드·설치·초기설정·로그인·점수입력")
    print("  추가 (14~22) — 타이머·일시정지·임시학생·AI코멘트·미리보기·카카오발송·완료")


if __name__ == "__main__":
    main()
