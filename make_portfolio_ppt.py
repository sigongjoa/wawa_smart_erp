#!/usr/bin/env python3
"""Wawa Smart ERP — Portfolio PPT Generator"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import os

SCREENSHOTS = "apps/desktop/e2e-screenshots-all"
OUTPUT = "Wawa_Smart_ERP_Portfolio.pptx"
GUIDE_OUTPUT = "Wawa_Smart_ERP_UserGuide.pptx"

# ── 색상 팔레트 ──────────────────────────────────────────
WAWA_BLUE   = RGBColor(0x1A, 0x73, 0xE8)   # 메인 블루
WAWA_NAVY   = RGBColor(0x0D, 0x1B, 0x2A)   # 다크 배경
KAKAO_YLW   = RGBColor(0xFF, 0xE8, 0x12)   # 카카오 노랑
ACCENT_ORG  = RGBColor(0xFF, 0x6B, 0x35)   # 강조 주황
ACCENT_GRN  = RGBColor(0x34, 0xA8, 0x53)   # 초록
TEXT_WH     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DK     = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_GY     = RGBColor(0x6B, 0x7C, 0x93)
BG_LIGHT    = RGBColor(0xF8, 0xF9, 0xFA)
BG_CARD     = RGBColor(0xEF, 0xF3, 0xFF)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=TEXT_DK,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb


def add_image_fit(slide, path, x, y, w, h):
    """이미지를 비율 유지하며 박스 안에 배치"""
    if not os.path.exists(path):
        add_rect(slide, x, y, w, h, fill=RGBColor(0xDD,0xDD,0xDD))
        add_text(slide, f"[missing]\n{os.path.basename(path)}", x, y, w, h,
                 size=9, color=TEXT_GY, align=PP_ALIGN.CENTER)
        return
    from PIL import Image
    img = Image.open(path)
    iw, ih = img.size
    ratio = min(w / (iw * 9144), h / (ih * 9144))  # EMU per px ~9144
    nw = int(iw * 9144 * ratio)
    nh = int(ih * 9144 * ratio)
    ox = x + (w - nw) // 2
    oy = y + (h - nh) // 2
    slide.shapes.add_picture(path, ox, oy, nw, nh)


def pill(slide, text, x, y, color, text_color=TEXT_WH, size=11):
    w, h = Cm(3.5), Cm(0.7)
    add_rect(slide, x, y, w, h, fill=color)
    add_text(slide, text, x, y, w, h, size=size, bold=True,
             color=text_color, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
#  SLIDES
# ════════════════════════════════════════════════════════

def slide_cover(prs):
    s = blank_slide(prs)
    fill_bg(s, WAWA_NAVY)

    # 좌측 색상 블록
    add_rect(s, Inches(0), Inches(0), Inches(0.6), SLIDE_H, fill=WAWA_BLUE)

    # 제목
    add_text(s, "Wawa Smart ERP", Inches(1), Inches(1.5), Inches(8), Inches(1.2),
             size=48, bold=True, color=TEXT_WH)
    add_text(s, "학원 운영의 시작부터 끝까지",
             Inches(1), Inches(2.9), Inches(8), Inches(0.7),
             size=22, color=RGBColor(0xA0,0xC4,0xFF))

    # 태그 라인
    add_text(s, "실시간 수업 타이머  ·  AI 월말 보고서  ·  카카오 1클릭 발송",
             Inches(1), Inches(3.8), Inches(10), Inches(0.5),
             size=14, color=RGBColor(0x80,0x9B,0xBD))

    # 하단 구분선 + 메타
    add_rect(s, Inches(1), Inches(6.5), Inches(4), Cm(0.12), fill=WAWA_BLUE)
    add_text(s, "Portfolio  ·  2026", Inches(1), Inches(6.7), Inches(4), Inches(0.4),
             size=11, color=RGBColor(0x60,0x75,0x8A))

    # 우측 장식 원
    add_rect(s, Inches(9.5), Inches(-1), Inches(5), Inches(5),
             fill=RGBColor(0x1E, 0x3A, 0x5F))
    add_rect(s, Inches(10.5), Inches(3.5), Inches(3.5), Inches(3.5),
             fill=RGBColor(0x16, 0x2B, 0x45))


def slide_problem(prs):
    s = blank_slide(prs)
    fill_bg(s, BG_LIGHT)

    add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(1.2), fill=WAWA_NAVY)
    add_text(s, "기존 학원의 불편함", Inches(0.5), Inches(0.2), Inches(8), Inches(0.8),
             size=28, bold=True, color=TEXT_WH)

    # 좌우 비교 카드
    col = [
        ("기존 방식 😓", ACCENT_ORG, [
            "수기 출석부 or 엑셀",
            "수업 연장 → 구두 메모",
            "점수 정리 → 카카오 수동 타이핑",
            "코멘트 작성 → 학생 1명당 30분",
            "카카오 개별 전송 → 30명 × 1시간",
        ]),
        ("Wawa Smart ERP ✅", ACCENT_GRN, [
            "실시간 타이머 — 자동 출결",
            "수업 연장 → +10/20/30분 버튼",
            "점수 입력 → 앱에서 바로 저장",
            "AI 코멘트 자동 생성 → 30초",
            "카카오 일괄 발송 → 1클릭",
        ]),
    ]
    for i, (title, color, items) in enumerate(col):
        cx = Inches(0.5 + i * 6.4)
        add_rect(s, cx, Inches(1.4), Inches(6), Inches(5.7),
                 fill=TEXT_WH, line=color, line_w=Pt(2))
        add_rect(s, cx, Inches(1.4), Inches(6), Inches(0.6), fill=color)
        add_text(s, title, cx + Inches(0.15), Inches(1.45), Inches(5.7), Inches(0.5),
                 size=15, bold=True, color=TEXT_WH)
        for j, item in enumerate(items):
            add_text(s, f"• {item}",
                     cx + Inches(0.2), Inches(2.2 + j * 0.88), Inches(5.6), Inches(0.8),
                     size=13, color=TEXT_DK)


def slide_section(prs, title, subtitle, color=WAWA_BLUE):
    s = blank_slide(prs)
    fill_bg(s, WAWA_NAVY)
    add_rect(s, Inches(0), Inches(3.0), SLIDE_W, Inches(0.08), fill=color)
    add_text(s, title, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
             size=42, bold=True, color=TEXT_WH, align=PP_ALIGN.CENTER)
    add_text(s, subtitle, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
             size=18, color=RGBColor(0xA0,0xC4,0xFF), align=PP_ALIGN.CENTER)


def slide_one_img(prs, title, note, img_path, badge=None, badge_color=WAWA_BLUE):
    s = blank_slide(prs)
    fill_bg(s, BG_LIGHT)

    # 상단 헤더
    add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(1.0), fill=WAWA_NAVY)
    add_text(s, title, Inches(0.4), Inches(0.1), Inches(9), Inches(0.8),
             size=22, bold=True, color=TEXT_WH)
    if badge:
        add_rect(s, Inches(10.5), Inches(0.2), Inches(2.3), Inches(0.6), fill=badge_color)
        add_text(s, badge, Inches(10.5), Inches(0.2), Inches(2.3), Inches(0.6),
                 size=12, bold=True, color=TEXT_WH, align=PP_ALIGN.CENTER)

    # 이미지
    add_image_fit(s, img_path, Inches(0.4), Inches(1.1), Inches(12.5), Inches(5.5))

    # 하단 노트
    add_rect(s, Inches(0), Inches(6.8), SLIDE_W, Inches(0.7),
             fill=RGBColor(0xE8,0xF0,0xFE))
    add_text(s, f"▶  {note}", Inches(0.3), Inches(6.82), Inches(12.5), Inches(0.5),
             size=12, color=WAWA_BLUE)


def slide_two_img(prs, title, img_left, img_right, cap_left, cap_right):
    s = blank_slide(prs)
    fill_bg(s, BG_LIGHT)

    add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(1.0), fill=WAWA_NAVY)
    add_text(s, title, Inches(0.4), Inches(0.1), Inches(12), Inches(0.8),
             size=22, bold=True, color=TEXT_WH)

    for i, (img, cap) in enumerate([(img_left, cap_left), (img_right, cap_right)]):
        x = Inches(0.3 + i * 6.6)
        add_image_fit(s, img, x, Inches(1.1), Inches(6.2), Inches(5.0))
        add_rect(s, x, Inches(6.2), Inches(6.2), Inches(0.6),
                 fill=RGBColor(0x1A,0x73,0xE8) if i==0 else ACCENT_ORG)
        add_text(s, cap, x, Inches(6.22), Inches(6.2), Inches(0.5),
                 size=11, bold=True, color=TEXT_WH, align=PP_ALIGN.CENTER)


def slide_three_img(prs, title, imgs, caps):
    """이미지 3장 가로 배열"""
    s = blank_slide(prs)
    fill_bg(s, BG_LIGHT)

    add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(1.0), fill=WAWA_NAVY)
    add_text(s, title, Inches(0.4), Inches(0.1), Inches(12), Inches(0.8),
             size=22, bold=True, color=TEXT_WH)

    colors = [ACCENT_GRN, ACCENT_ORG, RGBColor(0xE5,0x39,0x35)]
    for i in range(3):
        x = Inches(0.2 + i * 4.37)
        add_image_fit(s, imgs[i], x, Inches(1.1), Inches(4.1), Inches(4.9))
        add_rect(s, x, Inches(6.1), Inches(4.1), Inches(0.6), fill=colors[i])
        add_text(s, caps[i], x, Inches(6.12), Inches(4.1), Inches(0.5),
                 size=11, bold=True, color=TEXT_WH, align=PP_ALIGN.CENTER)


def slide_flow(prs):
    """리포트 전체 흐름 다이어그램 슬라이드"""
    s = blank_slide(prs)
    fill_bg(s, WAWA_NAVY)

    add_text(s, "월말 보고서 — 전체 흐름", Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
             size=28, bold=True, color=TEXT_WH)

    steps = [
        ("① 점수 입력", "선생님이 과목별\n점수 직접 입력", WAWA_BLUE),
        ("② AI 코멘트", "Gemini·Claude·GPT\n자동 생성 (30초)", ACCENT_ORG),
        ("③ 미리보기", "6개월 트렌드 차트\n과목별 코멘트 확인", ACCENT_GRN),
        ("④ 카카오 발송", "학부모에게 1클릭\n일괄 전송", KAKAO_YLW),
    ]
    for i, (title, desc, color) in enumerate(steps):
        x = Inches(0.5 + i * 3.1)
        add_rect(s, x, Inches(1.3), Inches(2.8), Inches(3.5), fill=color)
        add_text(s, title, x + Inches(0.1), Inches(1.5), Inches(2.6), Inches(0.6),
                 size=16, bold=True, color=WAWA_NAVY if color==KAKAO_YLW else TEXT_WH)
        add_text(s, desc, x + Inches(0.1), Inches(2.2), Inches(2.6), Inches(1.5),
                 size=13, color=WAWA_NAVY if color==KAKAO_YLW else TEXT_WH)

        if i < 3:
            add_text(s, "▶", Inches(3.2 + i * 3.1), Inches(2.4), Inches(0.5), Inches(0.5),
                     size=20, bold=True, color=TEXT_WH, align=PP_ALIGN.CENTER)

    # 하단 강조
    add_rect(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8),
             fill=RGBColor(0x1E,0x3A,0x5F))
    add_text(s,
             "기존: 선생님 1인 — 30명 × 5과목 코멘트 작성 → 약 3시간\n"
             "Wawa:  점수 입력 → AI 생성 → 발송 완료 → 약 10분",
             Inches(0.8), Inches(5.35), Inches(12), Inches(1.4),
             size=14, color=TEXT_WH)


def slide_moat(prs):
    s = blank_slide(prs)
    fill_bg(s, WAWA_NAVY)

    add_text(s, "왜 이 시스템이 경쟁 우위(Moat)를 갖는가",
             Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
             size=26, bold=True, color=TEXT_WH)

    moats = [
        ("🏃 운영 사이클 완전 커버",
         "타이머(입장) → 수업관리 → 월말평가 → 학부모 발송\n경쟁사는 이 중 1~2개만 제공",
         WAWA_BLUE),
        ("🤖 Multi-AI 추상화",
         "Gemini / Claude / GPT 3개 엔진 선택 가능\n비용 최적화 + 학원별 톤 커스텀",
         ACCENT_ORG),
        ("📱 카카오 알림톡 직연결",
         "학부모 열람률 95% 채널에 직접 발송\n이메일·문자 대비 압도적 도달률",
         KAKAO_YLW),
        ("🔒 오프라인 우선 Electron",
         "인터넷 불안정 지역에서도 동작\n학생 데이터 로컬 처리 → 개인정보보호",
         ACCENT_GRN),
        ("📊 Notion DB 연동",
         "학원이 기존에 쓰던 Notion 데이터 그대로 연결\n도입 마찰 최소화 + 이탈 비용 상승",
         RGBColor(0x6F,0x42,0xC1)),
        ("⏱ 실시간 타이머 차별화",
         "3단계 경고 (정상→주의→초과) + 정지 이유 기록\n수업 품질 분석 데이터 축적",
         RGBColor(0xE5,0x39,0x35)),
    ]
    for i, (title, desc, color) in enumerate(moats):
        row, col_i = divmod(i, 3)
        x = Inches(0.3 + col_i * 4.3)
        y = Inches(1.2 + row * 2.8)
        add_rect(s, x, y, Inches(4.1), Inches(2.5),
                 fill=RGBColor(0x1A,0x2E,0x44))
        add_rect(s, x, y, Inches(4.1), Inches(0.55), fill=color)
        add_text(s, title, x + Inches(0.1), y + Inches(0.05),
                 Inches(3.9), Inches(0.45),
                 size=13, bold=True,
                 color=WAWA_NAVY if color==KAKAO_YLW else TEXT_WH)
        add_text(s, desc, x + Inches(0.1), y + Inches(0.65),
                 Inches(3.9), Inches(1.7), size=11, color=RGBColor(0xC0,0xD4,0xEE))


# ════════════════════════════════════════════════════════
#  MAIN
# ════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    P = SCREENSHOTS + "/"

    # 1. 표지
    slide_cover(prs)

    # 2. 문제 정의
    slide_problem(prs)

    # ── TIMER SECTION ─────────────────────────────────
    slide_section(prs, "⏱ 실시간 수업 타이머",
                  "학생 입장부터 퇴원까지 — 한 화면으로 관리", WAWA_BLUE)

    slide_one_img(prs,
        "실시간 수업 현황 — 대기 학생 목록",
        "등원 예정 학생을 한눈에 확인 · 클릭 한 번으로 체크인",
        P + "rt_03_대기_학생_목록.png",
        badge="TIMER", badge_color=WAWA_BLUE)

    slide_one_img(prs,
        "체크인 후 — 수업 카드 실시간 타이머",
        "남은 시간 실시간 카운트다운 · 순수 수업시간 / 예정시간 동시 표시",
        P + "rt_04_체크인_수업카드.png",
        badge="REALTIME", badge_color=ACCENT_GRN)

    slide_two_img(prs,
        "일시정지 사유 선택  →  정지 상태 카드",
        P + "rt_05_일시정지_사유선택.png",
        P + "rt_06_일시정지_상태.png",
        "① 정지 버튼 → 사유 선택 (외출·휴식·화장실·기타)",
        "② 정지 중 카드 — 정지 시간 실시간 측정")

    slide_two_img(prs,
        "수업 연장 & 임시 학생 즉석 추가",
        P + "uc05-extend-sheet.png",
        P + "uc03-temp-modal.png",
        "+10 / +20 / +30분 또는 직접 입력으로 즉시 연장",
        "DB에 없는 학생도 수업 중 즉석 추가 — 임시 배지 표시")

    # ── REPORT SECTION ────────────────────────────────
    slide_section(prs, "📋 AI 월말 보고서",
                  "점수 입력 → AI 코멘트 → 카카오 발송 — 10분 완료", ACCENT_ORG)

    slide_flow(prs)

    slide_one_img(prs,
        "STEP 1 — 과목별 점수 입력",
        "선생님이 과목별 점수와 코멘트를 직접 입력 · 자동 저장",
        P + "08_성적표_성적_점수_입력.png",
        badge="STEP 1", badge_color=WAWA_BLUE)

    slide_one_img(prs,
        "STEP 2 — AI 코멘트 자동 생성",
        "Gemini · Claude · GPT 중 선택 → 버전별 코멘트 생성 → 마음에 드는 버전 적용",
        P + "13_성적표_AI_설정.png",
        badge="STEP 2", badge_color=ACCENT_ORG)

    slide_one_img(prs,
        "STEP 3 — 성적표 미리보기 (6개월 트렌드)",
        "최근 6개월 성적 트렌드 차트 · 과목별 코멘트 · 종합 평가 한 페이지로",
        P + "11_성적표_정지효_월말평가서.png",
        badge="STEP 3", badge_color=ACCENT_GRN)

    slide_one_img(prs,
        "STEP 4 — 카카오 알림톡 일괄 발송",
        "전송 대기 학생 일괄 선택 → 카카오 1클릭 발송 · 학부모 열람률 95%",
        P + "nav_08_step3_send.png",
        badge="STEP 4", badge_color=KAKAO_YLW)

    # ── MOAT ──────────────────────────────────────────
    slide_moat(prs)

    # 저장
    prs.save(OUTPUT)
    print(f"✅  저장 완료: {OUTPUT}")
    print(f"   슬라이드 수: {len(prs.slides)}")


if __name__ == "__main__":
    main()
