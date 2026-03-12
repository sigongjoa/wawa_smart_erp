#!/usr/bin/env python3
"""WAWA Smart ERP 사용자 가이드 PPT 생성 스크립트 v2"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── 경로 설정 ────────────────────────────────────────────────
BASE       = "/mnt/d/progress/wawa_smart_erp"
APP_SHOTS  = f"{BASE}/apps/desktop/ppt-screenshots"
ASSETS     = f"{BASE}/ppt-assets"
OUTPUT     = f"{BASE}/WAWA_SmartERP_사용자가이드.pptx"
GITHUB_URL = "github.com/sigongjoa/wawa_smart_erp/releases"

# ── 색상 팔레트 ──────────────────────────────────────────────
PRIMARY    = RGBColor(0x33, 0x72, 0xF7)
PRIMARY_DK = RGBColor(0x1A, 0x56, 0xDB)
ACCENT     = RGBColor(0xFF, 0x6B, 0x35)
RED        = RGBColor(0xDC, 0x26, 0x26)
RED_BG     = RGBColor(0xFF, 0xF0, 0xF0)
RED_DK     = RGBColor(0x7F, 0x1D, 0x1D)
GREEN      = RGBColor(0x10, 0xB9, 0x81)
GREEN_BG   = RGBColor(0xEC, 0xFD, 0xF5)
YELLOW     = RGBColor(0xFF, 0xD7, 0x00)
BG_LIGHT   = RGBColor(0xF0, 0xF4, 0xFF)
BG_PAGE    = RGBColor(0xF8, 0xFA, 0xFF)
BORDER     = RGBColor(0xE2, 0xE8, 0xF0)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x1E, 0x29, 0x3B)
MUTED      = RGBColor(0x64, 0x74, 0x8B)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── 헬퍼 ─────────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, sz=14, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb

def img(slide, path, l, t, w, h=None):
    if not os.path.exists(path):
        print(f"  ⚠ 이미지 없음: {path}")
        return
    kw = dict(width=Inches(w)) if not h else dict(width=Inches(w), height=Inches(h))
    slide.shapes.add_picture(path, Inches(l), Inches(t), **kw)

def header(slide, title, sub=""):
    rect(slide, 0, 0, 13.33, 1.2, fill=PRIMARY)
    txt(slide, title, 0.45, 0.1, 10, 0.72, sz=28, bold=True, color=WHITE)
    if sub:
        txt(slide, sub, 0.45, 0.77, 11.5, 0.38, sz=13,
            color=RGBColor(0xBF, 0xD7, 0xFF))

def badge(slide, n, l, t, color=ACCENT):
    rect(slide, l, t, 0.42, 0.42, fill=color)
    txt(slide, str(n), l+0.01, t-0.01, 0.40, 0.42,
        sz=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def info_box(slide, icon_text, body, l, t, w, h,
             bg=BG_LIGHT, border=PRIMARY, text_color=DARK):
    rect(slide, l, t, w, h, fill=bg, line=border, lw=Pt(1.8))
    txt(slide, icon_text, l+0.15, t+0.1, w-0.3, 0.35,
        sz=13, bold=True, color=border)
    txt(slide, body, l+0.15, t+0.45, w-0.3, h-0.55,
        sz=12, color=text_color, wrap=True)

def warn_box(slide, body, l, t, w, h):
    rect(slide, l, t, w, h, fill=RED_BG, line=RED, lw=Pt(2.5))
    txt(slide, "⚠️  주의", l+0.15, t+0.08, w-0.3, 0.38, sz=13, bold=True, color=RED)
    txt(slide, body, l+0.15, t+0.48, w-0.3, h-0.58, sz=12, color=RED_DK, wrap=True)

def ok_box(slide, body, l, t, w, h):
    rect(slide, l, t, w, h, fill=GREEN_BG, line=GREEN, lw=Pt(2))
    txt(slide, "✅  " + body, l+0.15, t+0.12, w-0.3, h-0.24,
        sz=13, bold=True, color=GREEN)

def screen_panel(slide, label, path, l=6.3, t=1.3, w=6.6, ph=5.7):
    rect(slide, l, t, w, ph, fill=BG_PAGE, line=BORDER, lw=Pt(1))
    txt(slide, label, l+0.15, t+0.1, w-0.3, 0.38,
        sz=12, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    if path and os.path.exists(path):
        img(slide, path, l+0.15, t+0.55, w-0.3)
    else:
        txt(slide, "(스크린샷 없음)", l+0.15, t+2.0, w-0.3, 1.0,
            sz=14, color=MUTED, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# 슬라이드 1: 표지
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=PRIMARY)
rect(slide, 0, 4.6, 13.33, 2.9, fill=PRIMARY_DK)

txt(slide, "WAWA Smart ERP", 1.0, 0.9, 11, 1.5,
    sz=64, bold=True, color=WHITE)
txt(slide, "사용자 가이드", 1.0, 2.55, 11, 0.9,
    sz=36, color=RGBColor(0xBF, 0xD7, 0xFF))

rect(slide, 1.0, 3.7, 9.0, 0.07, fill=RGBColor(0x93, 0xB8, 0xFF))

txt(slide, "다운로드 · 설치 · 시험 점수 입력까지\n처음 사용하는 선생님을 위한 단계별 안내",
    1.0, 3.9, 11, 1.0, sz=18, color=WHITE)
txt(slide, "v1.5.1  |  sigongjoa/wawa_smart_erp", 1.0, 6.6, 11, 0.5,
    sz=13, color=RGBColor(0x93, 0xB8, 0xFF))

# ════════════════════════════════════════════════════════════
# 슬라이드 2: 시작 전 준비물 (NEW)
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header(slide, "시작 전 준비물",
       "설치를 시작하기 전에 아래 두 가지를 관리자에게 먼저 받으세요")

txt(slide, "관리자(원장/원감)에게 미리 받아야 할 것", 0.5, 1.35, 12, 0.45,
    sz=15, bold=True, color=DARK)

# 준비물 1: config 파일
rect(slide, 0.5, 1.9, 5.8, 2.4, fill=BG_LIGHT, line=PRIMARY, lw=Pt(2))
rect(slide, 0.5, 1.9, 5.8, 0.5, fill=PRIMARY)
txt(slide, "① notion_config.json 파일", 0.7, 1.95, 5.4, 0.38,
    sz=15, bold=True, color=WHITE)
txt(slide, "앱과 Notion 데이터베이스를 연결하는 설정 파일입니다.",
    0.7, 2.52, 5.3, 0.45, sz=13, color=DARK)
txt(slide, "• 파일 이름: notion_config.json\n• 용도: 학생·성적 데이터 연결\n• 분실 시 관리자에게 재요청",
    0.7, 3.0, 5.3, 1.1, sz=12, color=MUTED)

# 준비물 2: PIN 번호
rect(slide, 6.9, 1.9, 5.8, 2.4, fill=BG_LIGHT, line=PRIMARY, lw=Pt(2))
rect(slide, 6.9, 1.9, 5.8, 0.5, fill=PRIMARY)
txt(slide, "② 본인 PIN 번호 (4자리)", 7.1, 1.95, 5.4, 0.38,
    sz=15, bold=True, color=WHITE)
txt(slide, "앱 로그인 시 사용하는 4자리 비밀번호입니다.",
    7.1, 2.52, 5.3, 0.45, sz=13, color=DARK)
txt(slide, "• 선생님마다 개별 번호 부여\n• 분실 시 관리자에게 연락\n• 타인에게 공유 금지",
    7.1, 3.0, 5.3, 1.1, sz=12, color=MUTED)

# 안내 박스
rect(slide, 0.5, 4.55, 12.2, 0.9, fill=RGBColor(0xFF, 0xF3, 0xE0),
     line=ACCENT, lw=Pt(1.5))
txt(slide, "💡  위 두 가지를 받지 못했다면 지금 바로 관리자에게 연락하세요. "
    "파일과 PIN이 없으면 앱을 사용할 수 없습니다.",
    0.7, 4.62, 11.8, 0.7, sz=13, color=DARK)

# 체크리스트
txt(slide, "준비 체크리스트", 0.5, 5.65, 5, 0.38, sz=13, bold=True, color=DARK)
txt(slide, "☐  notion_config.json 파일을 받았다\n☐  본인 PIN 번호를 받았다\n☐  Windows PC 준비 완료",
    0.5, 6.05, 5.5, 1.0, sz=13, color=DARK)

txt(slide, "준비가 완료됐으면\n다음 단계로 이동하세요 →", 8.5, 5.9, 4.5, 1.1,
    sz=15, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# 슬라이드 3: 다운로드
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header(slide, "① 다운로드", "아래 주소에서 설치 파일을 다운로드하세요 (Edge / Chrome 모두 가능)")

# URL 강조
rect(slide, 0.5, 1.35, 8.5, 0.75, fill=BG_LIGHT, line=PRIMARY, lw=Pt(2.5))
txt(slide, f"🔗  https://{GITHUB_URL}", 0.7, 1.43, 8.1, 0.55,
    sz=15, bold=True, color=PRIMARY)

steps_dl = [
    ("위 주소를 브라우저 주소창에 입력하고 Enter", ""),
    ("'Assets' 항목에서 WAWA-Smart-ERP-Setup-1.5.1.exe 클릭", ""),
    ("다운로드가 완료될 때까지 기다리기", "완료 전에 열면 오류 발생!"),
]
for i, (s, warn) in enumerate(steps_dl):
    badge(slide, i+1, 0.5, 2.3 + i*0.9)
    txt(slide, s, 1.1, 2.28 + i*0.9, 8.0, 0.45, sz=15, color=DARK)
    if warn:
        txt(slide, f"⚠ {warn}", 1.1, 2.7 + i*0.9, 8.0, 0.35, sz=12, color=RED)

# 오른쪽 스크린샷
screen_panel(slide, "다운로드 목록 화면 (Edge 기준)",
             f"{ASSETS}/smartscreen_01_download_list.png")

# 하단 팁
rect(slide, 0.5, 5.55, 5.5, 1.5, fill=BG_LIGHT, line=BORDER, lw=Pt(1))
txt(slide, "💡  Chrome 사용자", 0.7, 5.6, 5.0, 0.38, sz=13, bold=True, color=PRIMARY)
txt(slide, "Chrome에서는 우측 상단 다운로드 아이콘(↓)에서\n파일 목록을 확인할 수 있습니다.\nEdge와 경고 화면이 다를 수 있으나 방법은 동일합니다.",
    0.7, 5.98, 5.1, 1.0, sz=12, color=MUTED)

# ════════════════════════════════════════════════════════════
# 슬라이드 4: SmartScreen 경고 처리
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header(slide, "② 설치 — SmartScreen 경고 처리",
       "다운로드한 .exe를 열면 경고가 표시됩니다 — 아래 순서대로 진행하세요")

# 설명
rect(slide, 0.5, 1.32, 5.5, 0.85, fill=RGBColor(0xFF, 0xF3, 0xE0),
     line=ACCENT, lw=Pt(1.5))
txt(slide, "이 경고는 정상입니다!", 0.7, 1.37, 5.0, 0.38, sz=13, bold=True, color=ACCENT)
txt(slide, "코드 서명 인증서가 없는 앱에서 항상 표시됩니다.\nWAWA Smart ERP는 안전한 앱입니다.",
    0.7, 1.72, 5.0, 0.38, sz=12, color=DARK)

# 단계
steps_ss = [
    "다운로드된 .exe 파일을 클릭하여 열기",
    "경고 팝업에서  '...' (점 세 개) 버튼 클릭",
    "'그래도 계속' 선택",
]
for i, s in enumerate(steps_ss):
    badge(slide, i+1, 0.5, 2.4 + i*0.95)
    txt(slide, s, 1.1, 2.38 + i*0.95, 5.0, 0.65, sz=15, color=DARK)

ok_box(slide, "이후 설치 마법사가 실행됩니다!", 0.5, 5.45, 5.5, 0.7)

# 오른쪽: 경고 창
screen_panel(slide, "경고 창 — Edge 기준 (실제 화면)",
             f"{ASSETS}/smartscreen_03_warning.png", l=6.0, t=1.3, w=3.5, ph=5.7)

# '...' 클릭 후 메뉴
screen_panel(slide, "'...' 클릭 후 메뉴",
             f"{ASSETS}/smartscreen_02_context_menu.png", l=9.7, t=1.3, w=3.3, ph=5.7)

# ════════════════════════════════════════════════════════════
# 슬라이드 5: SmartScreen → 그래도 계속
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header(slide, "② 설치 — '그래도 계속' 선택", "마지막 확인 단계")

txt(slide, "'그래도 계속' 버튼이 있는 화면이 나타나면 클릭합니다.",
    0.5, 1.35, 8, 0.45, sz=15, color=DARK)

# 중앙 스크린샷
screen_panel(slide, "'그래도 계속' 화면 (실제 화면)",
             f"{ASSETS}/smartscreen_04_continue.png", l=2.0, t=1.85, w=5.5, ph=4.9)

ok_box(slide, "'그래도 계속' 클릭 → 설치 마법사가 열립니다",
       2.0, 6.9, 5.5, 0.55)

# 오른쪽 주의
rect(slide, 8.3, 1.85, 4.7, 4.9, fill=BG_LIGHT, line=BORDER, lw=Pt(1))
txt(slide, "버튼이 보이지 않는 경우", 8.5, 1.95, 4.3, 0.38,
    sz=13, bold=True, color=PRIMARY)
tips = [
    "① '...' 버튼을 먼저 클릭",
    "② 드롭다운에서 '그래도 계속' 선택",
    "③ 또는 '이 앱을 안전한 것으로 보고' 클릭 후 계속",
]
for i, t_ in enumerate(tips):
    txt(slide, t_, 8.5, 2.48 + i*0.6, 4.2, 0.5, sz=13, color=DARK)

warn_box(slide,
    "절대로 '삭제' 버튼을 누르지 마세요.\n파일이 삭제되어 다시 다운로드해야 합니다.",
    8.3, 4.7, 4.7, 0.95)

# ════════════════════════════════════════════════════════════
# 슬라이드 6: 설치 마법사 (NEW)
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header(slide, "② 설치 — 설치 마법사", "'그래도 계속' 후 설치 마법사를 따라 진행하세요")

steps_install = [
    ("Next", "설치 시작 — 'Next' 버튼 클릭"),
    ("Install", "'Install' 클릭 → 설치 진행 (1~2분 소요)"),
    ("Finish", "'Finish' 클릭 → 설치 완료, 앱 자동 실행"),
]
for i, (btn, desc) in enumerate(steps_install):
    y = 1.45 + i * 1.4
    rect(slide, 0.5, y, 1.1, 0.7, fill=PRIMARY)
    txt(slide, btn, 0.52, y+0.1, 1.06, 0.5,
        sz=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, desc, 1.75, y+0.12, 4.5, 0.5, sz=15, color=DARK)

    if i < 2:
        txt(slide, "↓", 0.9, y+0.78, 0.4, 0.45, sz=18, bold=True,
            color=MUTED, align=PP_ALIGN.CENTER)

ok_box(slide, "설치 완료! 바탕화면에 'WAWA Smart ERP' 아이콘이 생성됩니다.",
       0.5, 5.85, 5.5, 0.7)

# 오른쪽 설치 마법사 스크린샷
screen_panel(slide, "일반적인 설치 마법사 화면",
             f"{ASSETS}/installer_wizard.png", l=6.3, t=1.3, w=6.6, ph=5.7)

# ════════════════════════════════════════════════════════════
# 슬라이드 7: 초기 설정 (Notion 연동)
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header(slide, "③ 초기 설정", "앱 첫 실행 시 한 번만 하면 됩니다 — notion_config.json 업로드")

steps_setup = [
    "바탕화면의 'WAWA Smart ERP' 아이콘을 더블클릭",
    "아래와 같은 '초기 설정' 화면이 나타납니다",
    "'JSON 파일 업로드' 클릭 → 받아둔 notion_config.json 선택",
    "잠시 기다리면 로그인 화면으로 자동 이동합니다",
]
for i, s in enumerate(steps_setup):
    badge(slide, i+1, 0.5, 1.45 + i*0.9)
    txt(slide, s, 1.15, 1.43 + i*0.9, 5.2, 0.7, sz=14, color=DARK)

info_box(slide,
    "💡  이 설정은 한 번만!",
    "이후부터는 바로 로그인 화면으로 이동합니다.\n설정 파일을 분실하면 관리자에게 재요청하세요.",
    0.5, 5.15, 5.5, 1.25)

screen_panel(slide, "초기 설정 화면 (실제)",
             f"{APP_SHOTS}/01_초기설정_화면.png")

# ════════════════════════════════════════════════════════════
# 슬라이드 8: 로그인
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header(slide, "④ 로그인", "담당 선생님 계정으로 접속 — PIN 번호는 관리자에게 미리 받으세요")

login_steps = [
    ("선생님 선택",  "드롭다운 목록에서 본인 이름 선택"),
    ("PIN 입력",     "관리자에게 받은 4자리 PIN 번호 입력"),
    ("접속하기 클릭","메인 화면으로 이동"),
]
for i, (title, desc) in enumerate(login_steps):
    badge(slide, i+1, 0.5, 1.5 + i*1.1)
    txt(slide, title, 1.1, 1.48 + i*1.1, 5.0, 0.45, sz=16, bold=True, color=DARK)
    txt(slide, desc,  1.1, 1.92 + i*1.1, 5.0, 0.45, sz=13, color=MUTED)

# 로딩 안내
info_box(slide,
    "⏳  처음 로그인 시 로딩 시간 안내",
    "Notion 서버에서 데이터를 불러오므로\n최초 로그인 시 10~30초가 걸릴 수 있습니다.\n흰 화면이 보여도 기다려주세요.",
    0.5, 4.65, 5.5, 1.5)

# PIN 분실 안내
warn_box(slide, "PIN 번호를 잊어버린 경우\n관리자(원장/원감)에게 연락하세요.\n앱에서 직접 변경 불가", 0.5, 6.25, 5.5, 1.1)

screen_panel(slide, "로그인 화면 (실제)",
             f"{APP_SHOTS}/02_로그인_화면.png")

# ════════════════════════════════════════════════════════════
# 슬라이드 9: 메인 화면 소개
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header(slide, "메인 화면 소개", "로그인 완료 — 좌측 사이드바에서 메뉴를 선택하세요")

# 메인 스크린샷
img(slide, f"{APP_SHOTS}/03_메인_화면_로그인완료.png",
    0.3, 1.35, 9.0)

modules = [
    ("⏱ 시간표",   "학생 타이머 · 출석"),
    ("📊 성적표",   "시험 · 점수 · 리포트"),
    ("✏️ 채점",    "시험지 채점"),
    ("📋 보강관리", "결석 · 보강 일정"),
    ("👥 학생관리", "학생 정보 조회"),
    ("⚙️ 설정",    "시스템 설정"),
]
txt(slide, "주요 메뉴", 9.55, 1.35, 3.5, 0.4, sz=14, bold=True, color=PRIMARY)
for i, (name, desc) in enumerate(modules):
    rect(slide, 9.55, 1.82 + i*0.82, 3.45, 0.73, fill=BG_LIGHT, line=BORDER, lw=Pt(1))
    txt(slide, name, 9.72, 1.87 + i*0.82, 3.0, 0.36, sz=13, bold=True, color=DARK)
    txt(slide, desc, 9.72, 2.21 + i*0.82, 3.0, 0.28, sz=11, color=MUTED)

# ════════════════════════════════════════════════════════════
# 슬라이드 10: 시험 관리
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header(slide, "⑤ 시험 관리", "성적표 > 시험 관리 — 점수 입력 전에 먼저 시험 일정을 확인하세요")

# 경로
rect(slide, 0.5, 1.35, 5.5, 0.6, fill=BG_LIGHT, line=PRIMARY, lw=Pt(1.8))
txt(slide, "📊 성적표  ›  시험 관리", 0.7, 1.43, 5.0, 0.42, sz=15, bold=True, color=PRIMARY)

txt(slide, "현재 달 시험 현황", 0.5, 2.1, 5.5, 0.38, sz=14, bold=True, color=DARK)

# 자동 년월 안내
info_box(slide,
    "📅  년월은 자동으로 설정됩니다",
    "앱이 현재 날짜를 자동 인식해\n이번 달 데이터를 불러옵니다.\n별도로 설정할 필요가 없습니다.",
    0.5, 2.55, 5.5, 1.35)

exam_features = [
    "학생별 시험 일정 및 완료 여부 확인",
    "정기고사 날짜 일괄 설정",
    "시험지 템플릿 등록",
]
for i, f in enumerate(exam_features):
    badge(slide, i+1, 0.5, 4.15 + i*0.75, color=PRIMARY)
    txt(slide, f, 1.1, 4.13 + i*0.75, 5.0, 0.6, sz=14, color=DARK)

screen_panel(slide, "시험 관리 화면 (실제)",
             f"{APP_SHOTS}/05b_시험관리_일정탭.png")

# ════════════════════════════════════════════════════════════
# 슬라이드 11: 점수 입력 — 학생 선택
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header(slide, "⑥ 점수 입력 — 학생 선택",
       "성적표 > 성적 입력 — 왼쪽 학생 목록에서 학생을 선택하세요")

rect(slide, 0.5, 1.35, 5.5, 0.6, fill=BG_LIGHT, line=PRIMARY, lw=Pt(1.8))
txt(slide, "📊 성적표  ›  성적 입력", 0.7, 1.43, 5.0, 0.42, sz=15, bold=True, color=PRIMARY)

input_steps = [
    "왼쪽 학생 목록에서 이름 클릭",
    "이름 또는 초성으로 검색 가능 (예: 'ㅈㅈㅎ' → 정지효)",
    "선택하면 오른쪽에 과목별 입력 폼 표시",
]
for i, s in enumerate(input_steps):
    badge(slide, i+1, 0.5, 2.2 + i*0.9)
    txt(slide, s, 1.1, 2.18 + i*0.9, 5.0, 0.7, sz=14, color=DARK)

# 상태 배지 설명
txt(slide, "학생 목록의 상태 배지 의미", 0.5, 5.05, 5.5, 0.38,
    sz=13, bold=True, color=DARK)
badges = [("완료", GREEN, "모든 과목 점수 입력 완료"),
          ("진행중", RGBColor(0xD9, 0x77, 0x06), "일부 과목 입력됨"),
          ("미입력", MUTED, "아직 아무것도 입력 안됨")]
for i, (label, color, desc) in enumerate(badges):
    rect(slide, 0.5, 5.52 + i*0.55, 1.2, 0.38, fill=color)
    txt(slide, label, 0.52, 5.53 + i*0.55, 1.16, 0.34,
        sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, desc, 1.85, 5.55 + i*0.55, 3.9, 0.38, sz=12, color=DARK)

screen_panel(slide, "학생 목록 화면 (실제)",
             f"{APP_SHOTS}/06_점수입력_학생목록.png")

# ════════════════════════════════════════════════════════════
# 슬라이드 12: 점수 입력 — 정지효 예시 + 저장
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header(slide, "⑥ 점수 입력 — 점수 & 저장",
       "예시: 정지효 학생 선택 후 과목별 점수 입력 → 반드시 저장!")

txt(slide, "입력 항목", 0.5, 1.35, 5.5, 0.38, sz=14, bold=True, color=PRIMARY)

fields = [
    ("점수",     "0~100 숫자 입력"),
    ("코멘트",   "과목별 교사 한 마디 (선택 사항)"),
    ("종합평가", "학생 전체 종합 의견 (선택 사항)"),
]
for i, (title, desc) in enumerate(fields):
    y = 1.82 + i * 0.8
    rect(slide, 0.5, y, 5.5, 0.7, fill=WHITE, line=BORDER, lw=Pt(1))
    rect(slide, 0.5, y, 0.12, 0.7, fill=PRIMARY)
    txt(slide, title, 0.78, y+0.05, 2.0, 0.34, sz=14, bold=True, color=PRIMARY)
    txt(slide, desc,  0.78, y+0.38, 4.6, 0.27, sz=12, color=MUTED)

# 저장 경고 (강조)
rect(slide, 0.5, 4.27, 5.5, 1.3, fill=RED_BG, line=RED, lw=Pt(3))
txt(slide, "⚠️  반드시 '저장' 버튼을 눌러야 합니다!", 0.7, 4.32, 5.1, 0.45,
    sz=14, bold=True, color=RED)
txt(slide, "저장 버튼을 누르지 않으면 입력한 점수가\n모두 사라집니다. 과목별로 입력 후 즉시 저장하세요.",
    0.7, 4.77, 5.1, 0.7, sz=12, color=RED_DK)

# 학생 추가 문의
rect(slide, 0.5, 5.72, 5.5, 1.2, fill=BG_LIGHT, line=PRIMARY, lw=Pt(1.8))
txt(slide, "👤  학생이 목록에 없는 경우", 0.7, 5.78, 5.0, 0.38,
    sz=13, bold=True, color=PRIMARY)
txt(slide, "직접 추가 불가 — 관리자(원장/원감)에게\n학생 추가를 요청해 주세요.",
    0.7, 6.17, 5.0, 0.65, sz=12, color=MUTED)

# 오른쪽: 정지효 실제 스크린샷
screen_panel(slide, "점수 입력 화면 — 정지효 학생 (실제)",
             f"{APP_SHOTS}/07_점수입력_정지효_선택.png")

# ════════════════════════════════════════════════════════════
# 슬라이드 13: 저장 완료 확인 (NEW)
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header(slide, "⑥ 점수 입력 — 저장 완료 확인",
       "'저장' 버튼 클릭 후 이 화면이 나오면 성공입니다")

# 왼쪽 설명
txt(slide, "저장 성공 확인 방법", 0.5, 1.35, 5.5, 0.4,
    sz=14, bold=True, color=DARK)

ok_box(slide, "화면 상단에 초록색 알림이 표시되면 저장 완료!", 0.5, 1.85, 5.5, 0.7)

txt(slide, "저장 후 체크리스트", 0.5, 2.75, 5.5, 0.38, sz=13, bold=True, color=DARK)
checks = [
    "☑  초록 알림 '점수가 저장되었습니다' 확인",
    "☑  학생 목록 배지가 '완료' 또는 '진행중'으로 변경",
    "☑  다음 과목도 동일하게 입력 후 저장",
    "☑  모든 과목 완료 시 배지가 '완료'로 표시",
]
for i, c in enumerate(checks):
    txt(slide, c, 0.6, 3.2 + i*0.58, 5.2, 0.5, sz=13, color=DARK)

warn_box(slide,
    "저장 실패 시 빨간 알림이 표시됩니다.\n인터넷 연결을 확인하고 다시 시도하세요.",
    0.5, 5.55, 5.5, 1.0)

# 오른쪽: 토스트 알림 스크린샷
screen_panel(slide, "저장 완료 알림 (실제 화면)",
             f"{APP_SHOTS}/08b_점수입력_저장완료_토스트.png")

# ════════════════════════════════════════════════════════════
# 슬라이드 14: 완료 & 주의사항
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=PRIMARY)
rect(slide, 0, 4.5, 13.33, 3.0, fill=PRIMARY_DK)

txt(slide, "✅  기본 사용법 완료!", 0.8, 0.7, 11, 1.1,
    sz=52, bold=True, color=WHITE)
txt(slide, "점수 입력까지 성공적으로 마쳤습니다.", 0.8, 1.85, 11, 0.6,
    sz=20, color=RGBColor(0xBF, 0xD7, 0xFF))

rect(slide, 0.8, 2.6, 11, 0.06, fill=RGBColor(0x93, 0xB8, 0xFF))

txt(slide, "다음 단계로 해볼 수 있는 기능", 0.8, 2.75, 7, 0.45,
    sz=15, bold=True, color=RGBColor(0xBF, 0xD7, 0xFF))
next_steps = [
    "📊  성적표 미리보기 → 6개월 추세 차트 확인",
    "📱  카카오 알림톡 발송 → 학부모 성적 공유",
    "📋  보강관리 → 결석 처리 및 보강 일정 등록",
]
for i, s in enumerate(next_steps):
    txt(slide, s, 1.0, 3.3 + i*0.52, 7.2, 0.45, sz=15, color=WHITE)

# 오른쪽 주의사항 박스
rect(slide, 8.8, 2.6, 4.2, 4.7,
     fill=RGBColor(0x1A, 0x40, 0xAA), line=RGBColor(0x93, 0xB8, 0xFF), lw=Pt(1.5))
txt(slide, "꼭 기억하세요!", 9.0, 2.72, 3.8, 0.42,
    sz=14, bold=True, color=YELLOW)

notices = [
    ("⚠️  저장 필수",   "과목별 점수 입력 후\n반드시 '저장' 클릭"),
    ("👤  학생 추가",   "목록에 없는 학생 추가는\n관리자에게 문의"),
    ("🔑  PIN 분실",    "비밀번호 분실 시\n관리자에게 문의"),
    ("⏳  느린 로딩",   "10~30초 소요 정상,\n흰 화면도 기다리기"),
]
for i, (title, desc) in enumerate(notices):
    txt(slide, title, 9.0, 3.3 + i*1.0, 3.8, 0.38,
        sz=12, bold=True, color=YELLOW)
    txt(slide, desc, 9.0, 3.67 + i*1.0, 3.8, 0.5,
        sz=11, color=RGBColor(0xCC, 0xDD, 0xFF))

# ════════════════════════════════════════════════════════════
# 저장
# ════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"\n✅ PPT 생성 완료: {OUTPUT}")
print(f"   슬라이드 수: {len(prs.slides)}")
